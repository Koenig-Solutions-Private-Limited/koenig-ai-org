"""
Build slides.pptx for Cursor Composer 2 — Chapter 1: Composer 2 Models & IDE-First Workflow
13 slides, python-pptx, Koenig AI Academy branding.

Fixes applied vs. prior draft:
  - Title slide uses chapter name (not template placeholder text from body)
  - All bullet text word-wraps natively via pptx text frames; no char truncation
  - Every content slide has >= 3 substantive bullets
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---- Brand colors (Koenig AI Academy / Cursor-appropriate palette) ----
DARK_BG   = RGBColor(0x18, 0x18, 0x1B)   # zinc-900
LIGHT_BG  = RGBColor(0xF4, 0xF4, 0xF5)   # zinc-100
TEXT_LIGHT= RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x18, 0x18, 0x1B)
ACCENT    = RGBColor(0x63, 0x66, 0xF1)   # indigo-500
ACCENT2   = RGBColor(0x8B, 0x5C, 0xF6)   # violet-500
MUTED     = RGBColor(0xA1, 0xA1, 0xAA)   # zinc-400
HIGHLIGHT = RGBColor(0x6E, 0xE7, 0xB7)   # emerald-300 (for callouts)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank


def rect(slide, left, top, width, height, fill=None, line=None, line_pt=0):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.line.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    return s


def txtbox(slide, text, left, top, width, height,
           size=14, bold=False, color=TEXT_DARK,
           align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def header_bar(slide, bg=DARK_BG):
    rect(slide, 0, 0, W, H, bg)
    rect(slide, 0, 0, W, Inches(0.06), ACCENT)


def footer(slide):
    txtbox(slide, "academy.kspl.tech  |  Koenig AI Academy",
           Inches(0.5), H - Inches(0.38), Inches(12), Inches(0.3),
           size=9, color=MUTED)


def content_slide(title_text, bullets, bg=LIGHT_BG, accent=ACCENT):
    """Standard content slide with title + bullet list (3-5 bullets, full word-wrap)."""
    slide = prs.slides.add_slide(blank_layout)
    rect(slide, 0, 0, W, H, bg)
    rect(slide, 0, 0, W, Inches(0.06), accent)
    rect(slide, 0, Inches(0.06), Inches(0.25), H - Inches(0.06), accent)
    txtbox(slide, title_text,
           Inches(0.55), Inches(0.18), Inches(12.5), Inches(0.65),
           size=26, bold=True, color=TEXT_DARK if bg == LIGHT_BG else TEXT_LIGHT)
    rect(slide, Inches(0.55), Inches(0.88), Inches(12.5), Pt(1), fill=MUTED)
    y = Inches(1.05)
    for b in bullets:
        bullet_box = slide.shapes.add_textbox(Inches(0.55), y, Inches(12.3), Inches(1.1))
        tf = bullet_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"▸  {b}"
        run.font.size = Pt(15)
        run.font.color.rgb = TEXT_DARK if bg == LIGHT_BG else TEXT_LIGHT
        # advance y by approximate line height; pptx will wrap automatically
        y += Inches(0.95)
    footer(slide)
    return slide


# ---------------------------------------------------------------------------
# Slide 1: Title
# ---------------------------------------------------------------------------
def slide_01_title():
    slide = prs.slides.add_slide(blank_layout)
    rect(slide, 0, 0, W, H, DARK_BG)
    # Accent stripe
    stripe_w = W / 3
    for i, c in enumerate([ACCENT, ACCENT2, RGBColor(0x10, 0xB9, 0x81)]):
        rect(slide, stripe_w * i, 0, stripe_w, Inches(0.07), c)
    # Chapter label
    txtbox(slide, "CHAPTER 1", Inches(1), Inches(1.4), Inches(11.3), Inches(0.5),
           size=15, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
    # Chapter title — uses frontmatter title, NOT body template text
    txtbox(slide, "Composer 2 Models & IDE-First Workflow (2026)",
           Inches(1), Inches(2.0), Inches(11.3), Inches(1.6),
           size=38, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    # Subtitle
    txtbox(slide, "Cursor Composer 2  ·  Routing Rubric  ·  .cursorrules  ·  AGENTS.md",
           Inches(1), Inches(3.8), Inches(11.3), Inches(0.6),
           size=16, color=MUTED, align=PP_ALIGN.CENTER)
    # Course tag
    txtbox(slide, "Koenig AI Academy  |  academy.kspl.tech",
           Inches(1), Inches(5.1), Inches(11.3), Inches(0.4),
           size=12, color=MUTED, align=PP_ALIGN.CENTER)
    # Duration badge
    rect(slide, Inches(5.6), Inches(6.2), Inches(2.1), Inches(0.5), ACCENT)
    txtbox(slide, "⏱  45 min",
           Inches(5.65), Inches(6.22), Inches(2.0), Inches(0.42),
           size=13, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    notes(slide,
          "Chapter 1 grounds learners in what Cursor Composer 2 actually is — model lineage, "
          "pricing, honest benchmark reading — and gives them two foundational project files: "
          ".cursorrules and AGENTS.md. Every subsequent chapter assumes these are set up.")


# ---------------------------------------------------------------------------
# Slide 2: Learning Objectives
# ---------------------------------------------------------------------------
def slide_02_objectives():
    slide = prs.slides.add_slide(blank_layout)
    rect(slide, 0, 0, W, H, LIGHT_BG)
    rect(slide, 0, 0, W, Inches(0.06), ACCENT)
    txtbox(slide, "Learning Objectives",
           Inches(0.5), Inches(0.18), Inches(12), Inches(0.65),
           size=26, bold=True, color=TEXT_DARK)
    rect(slide, Inches(0.5), Inches(0.88), Inches(12.3), Pt(1), fill=MUTED)
    objs = [
        "Understand what Cursor Composer 2 is, including its model lineage and pricing",
        "Apply a routing rubric to decide when to use Composer 2 versus a frontier reasoning model",
        "Configure a project .cursorrules file that steers AI behavior reliably",
        "Set up an AGENTS.md file for multi-agent project context",
        "Enable informed benchmark reading with a two-layer skepticism framework",
    ]
    y = Inches(1.05)
    for i, obj in enumerate(objs):
        num_box = slide.shapes.add_textbox(Inches(0.5), y, Inches(0.45), Inches(0.75))
        tf_n = num_box.text_frame
        p_n = tf_n.paragraphs[0]
        p_n.alignment = PP_ALIGN.CENTER
        run_n = p_n.add_run()
        run_n.text = str(i + 1)
        run_n.font.size = Pt(18)
        run_n.font.bold = True
        run_n.font.color.rgb = ACCENT
        obj_box = slide.shapes.add_textbox(Inches(1.1), y + Inches(0.05), Inches(11.7), Inches(0.75))
        tf_o = obj_box.text_frame
        tf_o.word_wrap = True
        p_o = tf_o.paragraphs[0]
        run_o = p_o.add_run()
        run_o.text = obj
        run_o.font.size = Pt(15)
        run_o.font.color.rgb = TEXT_DARK
        y += Inches(1.05)
    footer(slide)
    notes(slide,
          "After this chapter, learners can explain the Composer 2 model lineage, apply the routing "
          "rubric to their own workloads, and configure the two project files that govern all AI "
          "sessions in a Cursor project.")


# ---------------------------------------------------------------------------
# Slide 3: Section 1.1 — What Composer 2 Is (model lineage + pricing)
# ---------------------------------------------------------------------------
def slide_03_what_is_composer2():
    s = content_slide(
        "1.1  What Cursor Composer 2 Is",
        [
            "Built on Kimi K2.5 base model via two-phase training: continued pretraining + "
            "large-scale reinforcement learning optimized for agentic IDE tasks "
            "(source: Cursor technical report, arXiv 2603.24477)",
            "Cursor did not prominently disclose the base model at launch — community analysis "
            "on HackerNews and r/LocalLLaMA identified the Kimi K2.5 lineage via architecture "
            "and training methodology comparison",
            "Model lineage matters in practice: base-model strengths and failure modes carry "
            "forward after fine-tuning; run your own evals, don't rely on vendor benchmarks alone",
            "Standard pricing: $0.50 / M input tokens, $2.50 / M output tokens — "
            "approximately 86 % cheaper than previous Composer generation (VentureBeat, 2026-05-14)",
            "Max plan: Composer 2 usage included in subscription, removing per-token cost "
            "friction for teams running 50+ agent turns per day",
        ]
    )
    notes(s,
          "Key teaching point: lineage disclosure matters not for drama but because it determines "
          "failure modes. Teams using Composer 2 for unfamiliar architecture work may hit the same "
          "gaps as Kimi K2.5. Price point changes iteration economics significantly.")


# ---------------------------------------------------------------------------
# Slide 4: Section 1.2 — Benchmark Literacy
# ---------------------------------------------------------------------------
def slide_04_benchmark_literacy():
    s = content_slide(
        "1.2  Benchmark Literacy: Reading Numbers Honestly",
        [
            "Cursor reports Composer 2 at 61.3 on CursorBench, positioned above Claude Opus 4.6; "
            "VentureBeat: 'beats Opus 4.6 but trails GPT-5.4' — useful directional signal only",
            "Layer 1 — Directional signal: a model scoring 61.3 vs 45 on the same eval likely "
            "handles those task types better in tested conditions; CursorBench tests IDE-specific "
            "agentic tasks (multi-file edits, test generation, tool use) — more relevant than "
            "generic coding competitions",
            "Layer 2 — Local task evals as decision authority: your codebase has specific language "
            "versions, tribal knowledge, test harnesses, and reviewer standards that no public "
            "benchmark captures; run Composer 2 on 10–20 representative tasks before committing",
            "Benchmaxxed risk: if a benchmark score rose sharply without proportionally large "
            "training compute, the model may be overfit to the eval distribution — "
            "treat like a suspiciously high A/B test and run your own experiment",
            "Measure what matters in your context: correctness rate, follow-up prompts needed, "
            "test-pass rate on first attempt, and token spend per task completion",
        ]
    )
    notes(s,
          "The two-layer framework prevents both dismissal ('benchmarks mean nothing') and "
          "credulity ('it's ranked higher so we'll always use it'). The local eval habit is "
          "the practical takeaway — it is the only way to know if the model suits your stack.")


# ---------------------------------------------------------------------------
# Slide 5: Section 1.3 — Routing Rubric: Use Composer 2 for…
# ---------------------------------------------------------------------------
def slide_05_routing_use_composer2():
    s = content_slide(
        "1.3  Routing Rubric — Use Composer 2 For:",
        [
            "Short-to-medium iteration loops: feature additions, bug fixes, refactor passes "
            "in code you own and understand well",
            "High-frequency experiments: parallel attempts ('try three approaches, I'll pick "
            "the best'), throwaway prototypes where iteration speed matters more than depth",
            "Routine generation tasks: CRUD endpoints, test cases, migration scripts for "
            "known schemas — tasks where you can predict what 'correct' looks like",
            "IDE-integrated edits: diff-reviewed changes where you inspect every hunk before "
            "accepting; Composer 2's cost economics support high-frequency review loops",
            "Cost-governed parallel agent sessions: token economics make Composer 2 the "
            "default for any task class where 50+ turns per day is feasible",
        ]
    )
    notes(s,
          "The routing rubric is the chapter's most practically useful artifact. "
          "Ask learners: which of their weekly tasks fall into these categories? "
          "Most iterative feature work does.")


# ---------------------------------------------------------------------------
# Slide 6: Section 1.3 — Routing Rubric: Escalate to Frontier
# ---------------------------------------------------------------------------
def slide_06_routing_escalate():
    s = content_slide(
        "1.3  Routing Rubric — Escalate to Frontier Model When:",
        [
            "Unfamiliar architecture: debugging third-party library internals you have never "
            "seen — when you cannot predict what the right answer looks like",
            "Brittle refactors: changing a core abstraction that 40+ files depend on, where "
            "the failure mode is subtle and a wrong answer may not be obviously wrong",
            "Security-sensitive reasoning: auth flows, crypto implementations, permission "
            "model design — areas where silent failure risk is unacceptable",
            "Complex cross-module debugging: tracing a bug that requires understanding state "
            "transitions across 5+ subsystems simultaneously",
            "Key escalation signal: 'When I cannot predict what correct looks like, and a "
            "wrong answer might not look wrong' — build this rubric into your .cursorrules "
            "as an explicit escalation instruction",
        ]
    )
    notes(s,
          "Tip: have learners add the escalation signal as a named rule in their .cursorrules: "
          "'If the task requires reasoning across 5+ unfamiliar subsystems, flag for human review "
          "before proceeding.' This makes the routing rubric machine-readable.")


# ---------------------------------------------------------------------------
# Slide 7: Section 1.4 — Auto Routing vs. Manual Pinning
# ---------------------------------------------------------------------------
def slide_07_auto_vs_pin():
    s = content_slide(
        "1.4  Auto Routing vs. Manual Model Pinning",
        [
            "Auto model selection delegates to Cursor's internal logic based on task type and "
            "context size — reasonable default for everyday edits and short sessions",
            "Pin manually for: any session running more than 10–15 turns (cost governance), "
            "parallel agent loops needing consistent behavior for comparison, and onboarding "
            "demos requiring reproducible output",
            "Pin manually for regulated environments: model selection must be auditable; "
            "Auto routing optimizes for Cursor's metrics, not yours",
            "To pin in Cursor 0.44: open Composer panel → model dropdown (top-right of input) "
            "→ select cursor-composer-2 → verify persistence in Settings → AI → Default Model",
            "Auto routing anti-pattern: using Auto for production-critical sessions or security-"
            "sensitive work where model identity must be known and stable",
        ]
    )
    notes(s,
          "The core message: Auto routing is a convenience default, not a production governance "
          "strategy. Teams that cannot name the model that generated a given change are not in "
          "a position to reason about its failure modes.")


# ---------------------------------------------------------------------------
# Slide 8: Section 1.5 — IDE Environment Setup
# ---------------------------------------------------------------------------
def slide_08_ide_setup():
    s = content_slide(
        "1.5  Setting Up the IDE Environment",
        [
            "Open the repository root (not a subdirectory): Composer uses the workspace root "
            "to discover .cursorrules and scope context-window crawls — opening a subdirectory "
            "silently breaks rule loading",
            "Enable inline diff mode (Settings → Editor → Diff: Inline Mode as of Cursor 0.44): "
            "makes AI-generated changes visibly separate from your code, reducing accidental "
            "acceptance of unwanted edits",
            "Use @file references explicitly rather than relying on open tabs being in context: "
            "Cursor includes open buffers in context — 20 open tabs is 20 files of noise",
            "Close files you are not actively editing: open buffers consume context window "
            "capacity; a disciplined tab hygiene habit is worth more than prompt engineering",
            "Start a new Composer session per discrete task: sessions accumulate context "
            "pollution; one new session per task is cheaper and more reliable than one "
            "mega-session per day",
        ]
    )
    notes(s,
          "These are the habits that separate teams that get consistent output from Composer 2 "
          "from those that fight it constantly. The workspace-root rule is the most commonly "
          "missed — many engineers open project subdirectories by habit from VS Code workflows.")


# ---------------------------------------------------------------------------
# Slide 9: Section 1.6 — Configuring .cursorrules
# ---------------------------------------------------------------------------
def slide_09_cursorrules():
    s = content_slide(
        "1.6  Configuring .cursorrules",
        [
            ".cursorrules is a project-level system prompt for Cursor: placed in the repository "
            "root, loaded at the start of every Composer session — the AI behavioral contract "
            "your whole team shares",
            "Weak rules fail because they are non-operationalizable: 'Write clean, readable "
            "code' gives the model nothing to verify against; 'Max function length: 40 lines, "
            "extract if longer' is verifiable",
            "Required categories in a production-ready file: stack pins (language, runtime, "
            "framework, DB, test runner, package manager), code standards, forbidden patterns, "
            "and explicit AI behavior instructions",
            "Forbidden patterns section is as important as standards: 'No raw SQL — use Drizzle "
            "query builder only', 'No console.log in production code — use lib/logger.ts', "
            "'No direct DOM manipulation outside React components'",
            "Verify your file is tracked: check git ls-files .cursorrules — some community "
            "starters add it to .gitignore, silently removing your team's shared AI contract",
        ]
    )
    notes(s,
          "The key principle: constraint-based rules outperform style-based rules. "
          "The model cannot operationalize vague style guidance, but it can check 'is this "
          "function longer than 40 lines?' Ask learners to list 3 verifiable rules for their "
          "own project before moving on.")


# ---------------------------------------------------------------------------
# Slide 10: Section 1.7 — Setting Up AGENTS.md
# ---------------------------------------------------------------------------
def slide_10_agents_md():
    s = content_slide(
        "1.7  Setting Up AGENTS.md",
        [
            "AGENTS.md serves a different purpose than .cursorrules: .cursorrules governs "
            "how the AI writes code; AGENTS.md provides project context — what the system "
            "is, how it is structured, what agents are active, and escalation rules",
            "Critical for multi-agent Cursor workflows: a Background Agent running a long "
            "task needs to understand the project without a live human in the loop; "
            "AGENTS.md is the document that grounds every agent session in real constraints",
            "Required sections: project purpose (one paragraph), directory architecture, "
            "active agents with scope limits ('no autonomous push rights — all changes "
            "require human review'), and ≥ 2 explicit escalation rules",
            "Escalation rules must name real risk boundaries: 'Any change to /server/auth "
            "→ flag for security review before commit', 'Any DB migration → run drizzle-kit "
            "check and confirm no destructive ops'",
            "Both files are needed together: .cursorrules alone leaves agents without "
            "project context; AGENTS.md alone leaves them without behavioral constraints",
        ]
    )
    notes(s,
          "AGENTS.md is the file that makes Cursor Background Agents safe to run on real "
          "codebases. Without it, the agent has no escalation path and will make "
          "architectural decisions silently. The escalation rules section is the "
          "highest-leverage part.")


# ---------------------------------------------------------------------------
# Slide 11: Section 1.8 — Common Pitfalls
# ---------------------------------------------------------------------------
def slide_11_pitfalls():
    s = content_slide(
        "1.8  Common Pitfalls & Anti-Patterns",
        [
            "Vague .cursorrules: 'Write clean code' is not enforceable; the model cannot "
            "operationalize 'clean' — replace every vague guideline with a specific, "
            "measurable constraint the model can verify",
            "One mega-session per day: running 60+ turns in a single session degrades output "
            "quality as context fills with earlier conversation turns — start a new session "
            "per discrete task",
            "Auto routing in production-critical sessions: for demos, security-sensitive work, "
            "or any session requiring reproducible behavior, pin the model explicitly; "
            "Auto optimizes for Cursor's metrics, not yours",
            "No AGENTS.md in multi-agent workflows: without project context and escalation "
            "rules, Background Agents have no boundary for what they must not do autonomously",
            "Treating benchmark superiority as task-specific superiority: Composer 2 outperforms "
            "some frontier models on CursorBench, but that benchmark does not represent your "
            "specific hard debugging tasks — always calibrate on your workload",
        ]
    )
    notes(s,
          "These five anti-patterns are distilled from community usage threads and direct "
          "production experience. The mega-session anti-pattern is the most universally "
          "underestimated — most developers do not realize context quality degrades "
          "significantly after 30+ turns.")


# ---------------------------------------------------------------------------
# Slide 12: Key Takeaways
# ---------------------------------------------------------------------------
def slide_12_takeaways():
    slide = prs.slides.add_slide(blank_layout)
    rect(slide, 0, 0, W, H, LIGHT_BG)
    rect(slide, 0, 0, W, Inches(0.06), ACCENT)
    txtbox(slide, "Chapter 1 — Key Takeaways",
           Inches(0.5), Inches(0.18), Inches(12), Inches(0.65),
           size=26, bold=True, color=TEXT_DARK)
    rect(slide, Inches(0.5), Inches(0.88), Inches(12.3), Pt(1), fill=MUTED)
    takeaways = [
        (ACCENT,  "Composer 2 = Kimi K2.5 + large-scale RL",
         "Two-phase training for agentic IDE tasks; $0.50/M input changes iteration economics"),
        (ACCENT2, "Two-layer benchmark reading",
         "Directional signal (rankings) + local task evals (your actual workload) — never one alone"),
        (RGBColor(0x10,0xB9,0x81), "Route by task class, not by preference",
         "Composer 2 for iteration loops; frontier models for unfamiliar architecture and security-"
         "sensitive reasoning"),
        (RGBColor(0xF5,0x9E,0x0B), ".cursorrules = constraint-based, not style-based",
         "Verifiable rules outperform vague guidelines; forbidden patterns section is as "
         "important as standards"),
        (RGBColor(0xEF,0x44,0x44), "AGENTS.md grounds every agent session",
         "Project context + escalation rules are the difference between safe and unsafe "
         "Background Agent runs"),
    ]
    y = Inches(1.02)
    for color, title, body in takeaways:
        rect(slide, Inches(0.5), y + Inches(0.1), Inches(0.18), Inches(0.8), color)
        title_box = slide.shapes.add_textbox(Inches(0.82), y, Inches(12), Inches(0.42))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        p_t = tf_t.paragraphs[0]
        run_t = p_t.add_run()
        run_t.text = title
        run_t.font.size = Pt(14)
        run_t.font.bold = True
        run_t.font.color.rgb = TEXT_DARK
        body_box = slide.shapes.add_textbox(Inches(0.82), y + Inches(0.4), Inches(12), Inches(0.55))
        tf_b = body_box.text_frame
        tf_b.word_wrap = True
        p_b = tf_b.paragraphs[0]
        run_b = p_b.add_run()
        run_b.text = body
        run_b.font.size = Pt(11)
        run_b.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
        y += Inches(1.1)
    footer(slide)
    notes(slide,
          "These five takeaways are the chapter's core — learners should be able to recite them "
          "before proceeding to Chapter 2. The most commonly forgotten: .cursorrules needs "
          "constraint-based (verifiable) rules, not style guidance.")


# ---------------------------------------------------------------------------
# Slide 13: What's Next
# ---------------------------------------------------------------------------
def slide_13_whats_next():
    slide = prs.slides.add_slide(blank_layout)
    rect(slide, 0, 0, W, H, DARK_BG)
    for i, c in enumerate([ACCENT, ACCENT2, RGBColor(0x10, 0xB9, 0x81)]):
        rect(slide, W * i / 3, H - Inches(0.07), W / 3, Inches(0.07), c)
    txtbox(slide, "Up Next",
           Inches(1), Inches(1.6), Inches(11.3), Inches(0.5),
           size=18, bold=True, color=ACCENT2, align=PP_ALIGN.CENTER)
    txtbox(slide, "Chapter 2: Mastering the Cursor Composer Interface",
           Inches(1), Inches(2.2), Inches(11.3), Inches(1.4),
           size=34, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    txtbox(slide,
           "Context window management as projects grow  ·  @file and @symbol reference precision  ·  "
           "Structuring multi-turn sessions that don't degrade",
           Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.9),
           size=15, color=MUTED, align=PP_ALIGN.CENTER)
    txtbox(slide,
           "Prerequisite: complete the hands-on exercise from Chapter 1 — "
           ".cursorrules and AGENTS.md must be committed before Chapter 2 workflows will work reliably",
           Inches(2), Inches(5.0), Inches(9.3), Inches(0.7),
           size=12, color=MUTED, align=PP_ALIGN.CENTER)
    notes(slide,
          "Chapter 2 builds directly on the project foundation from Chapter 1. "
          "The techniques for context window management only work well when .cursorrules "
          "and AGENTS.md are already in place. Make sure learners have completed the "
          "hands-on exercise before proceeding.")


# ---------------------------------------------------------------------------
# Build deck
# ---------------------------------------------------------------------------
slide_01_title()        # 1
slide_02_objectives()   # 2
slide_03_what_is_composer2()   # 3
slide_04_benchmark_literacy()  # 4
slide_05_routing_use_composer2()  # 5
slide_06_routing_escalate()    # 6
slide_07_auto_vs_pin()         # 7
slide_08_ide_setup()           # 8
slide_09_cursorrules()         # 9
slide_10_agents_md()           # 10
slide_11_pitfalls()            # 11
slide_12_takeaways()           # 12
slide_13_whats_next()          # 13

out = "vault/courses/cursor-composer-2/01-composer-models/slides.pptx"
prs.save(out)
print(f"OK  {out}  ({len(prs.slides)} slides)")
