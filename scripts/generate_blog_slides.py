#!/usr/bin/env python3
"""Generate slides.pptx for blog posts using python-pptx.

Brand: #18181b dark bg / #ffffff white text / #a1a1aa zinc-400 accent
Format: 5-7 slides (title + 3-5 content + CTA)
"""
import sys
sys.path.insert(0, '/tmp/pptx_deps')

import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

try:
    import yaml
    HAS_YAML = True
except ImportError:
    HAS_YAML = False

DARK   = RGBColor(0x18, 0x18, 0x1B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xA1, 0xA1, 0xAA)  # zinc-400

ATTRIBUTION = "academy.kspl.tech | Koenig AI Academy"

SKIP_SECTIONS = {
    "references", "further reading", "related content",
    "related articles", "notes", "acknowledgements",
}


def parse_frontmatter(text: str):
    if not text.startswith("---"):
        return {}, text
    end = text.find("---", 3)
    if end == -1:
        return {}, text
    fm_raw = text[3:end]
    body = text[end + 3:].strip()
    if HAS_YAML:
        try:
            fm = yaml.safe_load(fm_raw) or {}
        except Exception:
            fm = {}
    else:
        fm = {}
        for line in fm_raw.splitlines():
            if ":" in line:
                k, _, v = line.partition(":")
                fm[k.strip()] = v.strip()
    return fm, body


def clean_line(line: str) -> str:
    """Strip markdown syntax for display as plain bullet text."""
    line = re.sub(r'\*\*(.+?)\*\*', r'\1', line)
    line = re.sub(r'\*(.+?)\*', r'\1', line)
    line = re.sub(r'`([^`]+)`', r'\1', line)
    # Strip [[n]](url) footnote citations entirely (both ref and URL)
    line = re.sub(r'\[\[\d+\]\]\([^\)]+\)', '', line)
    line = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', line)
    # Obsidian wikilinks: [[path|display]] → display;
    # [[course/slug]] → academy URL; other [[path]] refs omitted
    line = re.sub(r'\[\[([^\]|]+)\|([^\]]+)\]\]', r'\2', line)
    def _wikilink(m):
        inner = m.group(1)
        if inner.startswith('course/'):
            return f"academy.kspl.tech/courses/{inner[len('course/'):]}"
        if inner.startswith('research/'):
            return 'the research note'
        return ''
    line = re.sub(r'\[\[([^\]]+)\]\]', _wikilink, line)
    line = re.sub(r'\s{2,}', ' ', line)  # collapse spaces from stripped refs
    line = re.sub(r'<[^>]+>', '', line)
    line = re.sub(r'^\s*[-*+]\s+', '', line)  # strip list markers
    line = re.sub(r'^\s*\d+\.\s+', '', line)  # strip numbered list
    line = line.strip()
    return line


def extract_sections(body: str):
    """Return list of (slide_title, [bullet_str, ...]).

    H2 sections with ≥4 H3 subsections are split into two slides so no
    single slide covers more than half the subsections.
    """
    # First pass: build H2 → H3 hierarchy
    raw: list[tuple] = []   # (h2_title, preamble_bullets, [(h3_title, [bullets])])

    current_h2 = None
    current_h3 = None
    cur_h3_bullets = []
    cur_h3_list = []
    preamble = []

    def _keep(cleaned):
        return (cleaned
                and len(cleaned) > 15
                and not cleaned.startswith("|")
                and not cleaned.startswith("```")
                and not cleaned.startswith("!["))

    def flush_h3():
        nonlocal current_h3, cur_h3_bullets
        if current_h3 is not None:
            cur_h3_list.append((current_h3, list(cur_h3_bullets[:4])))
            current_h3 = None
            cur_h3_bullets = []

    def flush_h2():
        nonlocal current_h2
        flush_h3()
        raw.append((current_h2, list(preamble), list(cur_h3_list)))
        current_h2 = None
        del cur_h3_list[:]
        del preamble[:]

    for raw_line in body.splitlines():
        if raw_line.startswith("## "):
            flush_h2()
            current_h2 = raw_line[3:].strip()
        elif raw_line.startswith("### ") and current_h2 is not None:
            flush_h3()
            current_h3 = raw_line[4:].strip()
        elif not raw_line.startswith("#") and current_h2 is not None:
            cleaned = clean_line(raw_line)
            if _keep(cleaned):
                if current_h3 is not None:
                    cur_h3_bullets.append(cleaned[:120])
                else:
                    preamble.append(cleaned[:120])
    flush_h2()

    # Second pass: convert to slides
    slides = []
    for h2_title, preamble_bullets, h3s in raw:
        if h2_title is None:
            continue  # pre-H2 preamble handled via frontmatter in generate()
        if h2_title.lower() in SKIP_SECTIONS:
            continue

        if len(h3s) >= 4:
            # Split dense H2: pair adjacent H3s into two slides
            mid = (len(h3s) + 1) // 2
            for group in (h3s[:mid], h3s[mid:]):
                grp_bullets = []
                for h3t, h3b in group:
                    short = re.sub(r'^\d+\.\s+', '', h3t).split(':')[0].strip()
                    grp_bullets.append(short[:80])
                    grp_bullets.extend(h3b[:2])
                if grp_bullets:
                    parts = [re.sub(r'^\d+\.\s+', '', h3t).split(':')[0].strip()[:40]
                             for h3t, _ in group]
                    slides.append((' & '.join(parts), grp_bullets[:5]))
        elif h3s:
            all_bullets = list(preamble_bullets[:2])
            for h3t, _ in h3s[:5]:
                all_bullets.append(re.sub(r'^\d+\.\s+', '', h3t)[:80])
            if all_bullets:
                slides.append((h2_title, all_bullets[:5]))
        elif preamble_bullets:
            slides.append((h2_title, preamble_bullets[:5]))

    return slides


def _add_text(slide, text, left, top, width, height,
              size, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK


def _add_footer(slide):
    _add_text(slide, ATTRIBUTION,
              Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.4),
              size=10, color=ACCENT)


def make_title_slide(prs, title: str, date_str: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, title,
              Inches(1.0), Inches(2.0), Inches(11.3), Inches(2.5),
              size=34, bold=True, align=PP_ALIGN.CENTER)
    if date_str:
        _add_text(slide, str(date_str),
                  Inches(1.0), Inches(4.7), Inches(11.3), Inches(0.5),
                  size=16, color=ACCENT, align=PP_ALIGN.CENTER)
    _add_footer(slide)


def make_content_slide(prs, section_title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, section_title,
              Inches(0.8), Inches(0.45), Inches(11.7), Inches(1.1),
              size=28, bold=True)
    # Horizontal rule via thin rectangle
    rect = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0.8), Inches(1.55), Inches(11.7), Inches(0.04)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = ACCENT
    rect.line.fill.background()

    top = Inches(1.75)
    for bullet in bullets[:4]:
        _add_text(slide, f"  {bullet}",
                  Inches(0.8), top, Inches(11.7), Inches(1.05),
                  size=17, color=WHITE)
        top += Inches(1.15)

    _add_footer(slide)


def make_cta_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, "Read the full article",
              Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.0),
              size=30, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "academy.kspl.tech",
              Inches(1.0), Inches(3.5), Inches(11.3), Inches(0.9),
              size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    _add_text(slide, "Koenig AI Academy",
              Inches(1.0), Inches(6.85), Inches(11.3), Inches(0.4),
              size=12, color=ACCENT, align=PP_ALIGN.CENTER)


def generate(draft_path: Path) -> Path:
    text = draft_path.read_text(encoding="utf-8")
    fm, body = parse_frontmatter(text)

    # Get title from frontmatter or first H1
    title = fm.get("title") or fm.get("h1") or ""
    if not title:
        h1 = re.search(r'^# (.+)', body, re.MULTILINE)
        title = h1.group(1).strip() if h1 else draft_path.parent.name.replace("-", " ").title()

    date_str = fm.get("date", "")

    sections = extract_sections(body)

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs, title, date_str)

    content_count = 0

    # Add "What You'll Learn" slide from learning_objectives when content is sparse
    objectives = fm.get('learning_objectives') or []
    if isinstance(objectives, list) and objectives and len(sections) < 5:
        obj_bullets = [str(o)[:100] for o in objectives[:5]]
        make_content_slide(prs, "What You'll Learn", obj_bullets)
        content_count += 1

    for sec_title, bullets in sections:
        if content_count >= 5:
            break
        make_content_slide(prs, sec_title, bullets)
        content_count += 1

    make_cta_slide(prs)

    out = draft_path.parent / "slides.pptx"
    prs.save(str(out))
    return out


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: generate_blog_slides.py <draft.md> [...]")
        sys.exit(1)

    ok = []
    fail = []
    for arg in sys.argv[1:]:
        p = Path(arg)
        if not p.exists():
            print(f"SKIP (not found): {p}")
            fail.append(str(p))
            continue
        try:
            out = generate(p)
            slide_count = len(Presentation(str(out)).slides)
            size_kb = out.stat().st_size // 1024
            print(f"OK  {out}  ({slide_count} slides, {size_kb} KB)")
            ok.append(str(out))
        except Exception as e:
            print(f"FAIL {p}: {e}")
            fail.append(str(p))

    print(f"\nDone: {len(ok)} OK, {len(fail)} failed")
