"""Build the canonical Chapter 4 slide deck for Gemini Enterprise Agents.

Chapter: Comparing to Claude Agent SDK + Cloudflare Agents
7 slides: title + 5 content + CTA.
Fixes all KOEA-8122 G0 blockers:
  - MAX_BULLET_CHARS=200 (all bullets <120 chars; word-wrap handles the rest)
  - Wikilinks resolved to display text (no [[...]] in any slide)
  - No Python code artifacts (content is curated, not extracted)
  - >=3 bullets per content slide
"""
import sys
sys.path.insert(0, '/tmp/pptx_deps')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Brand colours
DARK   = RGBColor(0x18, 0x18, 0x1B)   # #18181b dark background
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)   # #ffffff text
ACCENT = RGBColor(0xA1, 0xA1, 0xAA)  # #a1a1aa zinc-400
INDIGO = RGBColor(0x63, 0x66, 0xF1)  # #6366f1 highlight / CTA

W = Inches(13.333)
H = Inches(7.5)
ATTRIBUTION = "academy.kspl.tech | Koenig AI Academy"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]  # fully blank layout


# ── helpers ───────────────────────────────────────────────────────────────────

def _bg(slide):
    s = slide.shapes.add_shape(1, 0, 0, W, H)
    s.fill.solid()
    s.fill.fore_color.rgb = DARK
    s.line.fill.background()


def _rect(slide, left, top, width, height, fill):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def _text(slide, text, left, top, width, height,
          size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _footer(slide):
    _text(slide, ATTRIBUTION,
          Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
          size=10, color=ACCENT)


def _rule(slide, y=Inches(1.3)):
    r = slide.shapes.add_shape(1, Inches(0.5), y, Inches(12.3), Inches(0.03))
    r.fill.solid()
    r.fill.fore_color.rgb = ACCENT
    r.line.fill.background()


def content_slide(title: str, bullets: list):
    """Add a branded content slide with title + bullets."""
    slide = prs.slides.add_slide(blank)
    _bg(slide)
    _rect(slide, 0, 0, W, Inches(0.06), INDIGO)           # top accent bar
    _text(slide, title,
          Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.95),
          size=28, bold=True)
    _rule(slide, Inches(1.25))

    n = len(bullets)
    step  = Inches(0.90) if n >= 5 else Inches(1.05)
    box_h = Inches(0.80) if n >= 5 else Inches(0.95)
    fsize = 15 if n >= 5 else 17
    y = Inches(1.42)
    for bullet in bullets[:5]:
        _text(slide, f"▸  {bullet}",
              Inches(0.55), y, Inches(12.2), box_h,
              size=fsize, color=WHITE)
        y += step

    _footer(slide)
    return slide


# ── Slide 1: Title ────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank)
_bg(s1)
_rect(s1, 0, 0, W, Inches(0.06), INDIGO)
_text(s1, "CHAPTER 4",
      Inches(1), Inches(1.8), Inches(11.3), Inches(0.5),
      size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
_text(s1, "Comparing to Claude Agent SDK\n+ Cloudflare Agents",
      Inches(1), Inches(2.5), Inches(11.3), Inches(2.0),
      size=38, bold=True, align=PP_ALIGN.CENTER)
_text(s1, "Gemini Enterprise Agents — Koenig AI Academy",
      Inches(1), Inches(5.0), Inches(11.3), Inches(0.5),
      size=14, color=ACCENT, align=PP_ALIGN.CENTER)
_footer(s1)

# ── Slide 2: Learning Objectives ──────────────────────────────────────────────
content_slide(
    "Learning Objectives",
    [
        "Contrast GEAP state management with Claude Agent SDK and Cloudflare Durable Objects",
        "Identify the deployment topology differences across all three platforms",
        "Name three workloads where GEAP wins and three where a lighter alternative is preferable",
        "Apply a vendor-selection framework to a real-world scenario",
    ]
)

# ── Slide 3: Platform Overview ────────────────────────────────────────────────
content_slide(
    "Platform Overview — Three Approaches",
    [
        "GEAP: fully managed and opinionated — state via Agent Sessions + Memory Bank, runs on GCP",
        "Claude Agent SDK: code-first, least opinionated — you own infra, state, and deployment entirely",
        "Cloudflare Agents: TypeScript SDK on Workers — built-in SQLite state, 300+ edge PoPs, sub-ms cold starts",
        "All three support tool calling, multi-agent patterns, and long-running agents",
        "Key divergence: state management is the sharpest architectural difference between platforms",
    ]
)

# ── Slide 4: State Management ─────────────────────────────────────────────────
content_slide(
    "State Management — Deepest Divergence",
    [
        "GEAP: platform manages state via Agent Sessions and Memory Bank — no schema work, but GCP-locked",
        "Claude SDK: you manage everything — conversation history, long-term memory, and context compression",
        "Cloudflare: this.setState() on a Durable Object — atomic, co-located with compute, but proprietary",
        "GEAP wins on convenience; Claude SDK wins on portability; Cloudflare wins on co-located simplicity",
        "Memory Bank has no export API at launch — factor migration cost in before adopting GEAP for state",
    ]
)

# ── Slide 5: Deployment Topology ──────────────────────────────────────────────
content_slide(
    "Deployment Topology",
    [
        "GEAP: GCP regional (us-central1, europe-west4) — compliance features, GCP ecosystem integration",
        "Claude SDK: deploy on any infra — any cloud, on-premises, or hybrid; no platform lock",
        "Cloudflare: 300+ global edge PoPs, sub-millisecond cold starts, WebSocket native via Durable Objects",
        "Real-time workloads: Cloudflare wins on latency; regulated workloads: GEAP wins on governance",
        "Scheduling: GEAP via Cloud Scheduler, Cloudflare via Durable Object alarms, Claude SDK: bring your own",
    ]
)

# ── Slide 6: Decision Framework ───────────────────────────────────────────────
content_slide(
    "Decision Framework — Which Platform When",
    [
        "Choose GEAP when on GCP, needing enterprise governance, or orchestrating five or more agents",
        "Choose Claude SDK when reasoning quality is the priority or vendor lock-in must be minimised",
        "Choose Cloudflare Agents when latency is critical, the team is TypeScript-native, or state is simple",
        "Hybrid patterns work: GEAP orchestration + Claude Opus sub-agents; Cloudflare edge + GEAP backend",
    ]
)

# ── Slide 7: CTA ──────────────────────────────────────────────────────────────
s7 = prs.slides.add_slide(blank)
_bg(s7)
_rect(s7, 0, 0, Inches(0.06), H, INDIGO)
_text(s7, "Try It Next",
      Inches(1), Inches(2.0), Inches(11.3), Inches(0.8),
      size=36, bold=True, align=PP_ALIGN.CENTER)
_text(s7, "Capstone: Build a Two-Agent Invoice-Processing Pipeline",
      Inches(1), Inches(3.0), Inches(11.3), Inches(0.7),
      size=22, color=ACCENT, align=PP_ALIGN.CENTER)
_text(s7, ATTRIBUTION,
      Inches(1), Inches(4.2), Inches(11.3), Inches(0.5),
      size=16, color=INDIGO, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
OUT = "vault/courses/gemini-enterprise-agents/04-comparing-to-claude-agent-sdk-and-cloudflare-agents-slides.pptx"
prs.save(OUT)
slide_count = len(prs.slides)
size_kb = __import__("os").path.getsize(OUT) // 1024
print(f"OK  {OUT}  ({slide_count} slides, {size_kb} KB)")
