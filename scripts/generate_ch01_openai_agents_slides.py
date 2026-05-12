#!/usr/bin/env python3
"""Generate ch01-slides.pptx for openai-agents-sdk-mastery.

Chapter 1: Anatomy of an Agent [60 mins]
Matches the visual style of picking-a-frontier-model decks:
- 13.33×7.5 in slides
- Brand: #18181b dark bg / #ffffff white / #a1a1aa zinc-400 accent / #3b82f6 blue highlight
- Structure: title · agenda · objectives · concept slides · exercise · knowledge check · summary · CTA
"""
import sys
sys.path.insert(0, '/tmp/pptx_deps')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK   = RGBColor(0x18, 0x18, 0x1B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xA1, 0xA1, 0xAA)   # zinc-400
BLUE   = RGBColor(0x3B, 0x82, 0xF6)   # blue-500
GREEN  = RGBColor(0x22, 0xC5, 0x5E)   # green-500
YELLOW = RGBColor(0xFA, 0xCC, 0x15)   # yellow-400

COURSE = "OpenAI Agents SDK Mastery"
FOOTER = f"{COURSE}  ·  Koenig AI Academy"


# ── helpers ───────────────────────────────────────────────────────────────────

def _set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK


def _add_text(slide, text, left, top, width, height,
              size, bold=False, color=WHITE,
              align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _add_footer(slide, text=FOOTER):
    _add_text(slide, text,
              Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4),
              size=10, color=ACCENT)


def _add_rule(slide, top=Inches(1.55)):
    rect = slide.shapes.add_shape(1, Inches(0.8), top, Inches(11.7), Inches(0.04))
    rect.fill.solid()
    rect.fill.fore_color.rgb = ACCENT
    rect.line.fill.background()


def _section_header(slide, title, subtitle=None):
    _add_text(slide, title,
              Inches(0.8), Inches(0.35), Inches(11.7), Inches(1.1),
              size=28, bold=True)
    _add_rule(slide)
    if subtitle:
        _add_text(slide, subtitle,
                  Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.6),
                  size=14, color=ACCENT, italic=True)


def _bullets(slide, items, start_top=None, size=17, indent="▸  ", color=WHITE):
    top = start_top if start_top is not None else Inches(2.3)
    for item in items:
        _add_text(slide, f"{indent}{item}",
                  Inches(0.8), top, Inches(11.7), Inches(1.0),
                  size=size, color=color)
        top += Inches(0.92)


# ── slide builders ────────────────────────────────────────────────────────────

def slide_01_title(prs):
    """Chapter header — full-bleed dark."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _add_text(s, "OPENAI AGENTS SDK MASTERY",
              Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.6),
              size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    _add_text(s, "Chapter 1",
              Inches(0.8), Inches(1.85), Inches(11.7), Inches(0.7),
              size=22, bold=False, color=ACCENT, align=PP_ALIGN.CENTER)
    _add_text(s, "Anatomy of an Agent",
              Inches(0.8), Inches(2.7), Inches(11.7), Inches(2.2),
              size=40, bold=True, align=PP_ALIGN.CENTER)
    _add_text(s, "60 min  ·  Agents vs. chatbots · Agent loops · Tool definition · Runtime setup",
              Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.6),
              size=14, color=ACCENT, align=PP_ALIGN.CENTER)
    _add_text(s, "academy.kspl.tech | Koenig AI Academy",
              Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.4),
              size=10, color=ACCENT, align=PP_ALIGN.CENTER)


def slide_02_agenda(prs):
    """Course agenda — all 10 chapters, Chapter 1 highlighted."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _section_header(s, "Course Agenda", "10 chapters · 8 hours total")

    chapters = [
        ("Ch 1", "Anatomy of an Agent", True),
        ("Ch 2", "The Agent SDK Tool-Calling Model", False),
        ("Ch 3", "State and Persistent Memory", False),
        ("Ch 4", "Planning and Reasoning", False),
        ("Ch 5", "Human-in-the-loop Workflows", False),
        ("Ch 6", "Observability, Tracing, and Debugging", False),
        ("Ch 7", "Testing and Evaluation", False),
        ("Ch 8", "Production Deployment", False),
        ("Ch 9", "Advanced Patterns", False),
        ("Ch 10", "Capstone: Autonomous Research Assistant", False),
    ]

    top = Inches(2.1)
    col1_x = Inches(0.8)
    col2_x = Inches(1.8)

    for i, (label, title, current) in enumerate(chapters):
        color = BLUE if current else WHITE
        label_color = BLUE if current else ACCENT
        _add_text(s, label, col1_x, top, Inches(0.9), Inches(0.42),
                  size=12, bold=current, color=label_color)
        _add_text(s, ("→  " if current else "    ") + title,
                  col2_x, top, Inches(10.5), Inches(0.42),
                  size=12, bold=current, color=color)
        top += Inches(0.435)

    _add_footer(s)


def slide_03_objectives(prs):
    """Learning objectives for Chapter 1."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _section_header(s, "What you'll be able to do", "Learning Objectives — Chapter 1")
    objectives = [
        "Define agents vs. chatbots — what makes a system 'agentic'",
        "Identify the core SDK objects: Agent, Runner, Tool, and Thread",
        "Explain the agent loop — perceive, reason, act, observe, repeat",
        "Set up the OpenAI Agents SDK development environment from scratch",
        "Build and run a working 'Hello World' agent end-to-end",
    ]
    _bullets(s, objectives, start_top=Inches(2.3), size=17)
    _add_footer(s)


def slide_04_agents_vs_chatbots(prs):
    """Concept: Agents vs. chatbots."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _section_header(s, "Agents vs. Chatbots",
                    "The key architectural difference — and why it matters")

    rows = [
        ("Turns",     "Single-turn (one Q → one A)",          "Multi-turn loop until goal is met"),
        ("Initiative","Reactive — waits to be asked",          "Proactive — decides next action itself"),
        ("State",     "Stateless by default",                  "Maintains state across steps (Thread)"),
        ("Actions",   "Text response only",                    "Calls tools, writes files, calls APIs"),
        ("Goal",      "Answer the question",                   "Complete a task"),
    ]
    headers = ["Dimension", "Chatbot", "Agent"]
    widths  = [2.0, 4.5, 4.5]
    left_base = 0.8
    top = Inches(1.85)

    left = left_base
    for h, w in zip(headers, widths):
        _add_text(s, h, Inches(left), top, Inches(w), Inches(0.45),
                  size=13, bold=True, color=ACCENT)
        left += w
    top += Inches(0.5)

    for dim, chatbot, agent in rows:
        left = left_base
        for cell, w in zip([dim, chatbot, agent], widths):
            color = BLUE if cell == agent and cell != dim else WHITE
            color = ACCENT if cell == dim else color
            _add_text(s, cell, Inches(left), top, Inches(w), Inches(0.55),
                      size=13, color=color)
            left += w
        top += Inches(0.62)

    _add_text(s,
              "Key insight: an agent replaces a human operator, not just a search box.",
              Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.45),
              size=13, bold=True, color=BLUE)
    _add_footer(s)


def slide_05_agent_loop(prs):
    """Concept: The agent loop."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _section_header(s, "The Agent Loop",
                    "Core execution pattern — the heartbeat of every agent")

    steps = [
        ("Receive",  "Task or user message enters the agent's context (Thread)"),
        ("Reason",   "LLM call with system prompt + tools + current context → decides what to do next"),
        ("Act",      "If a tool call is returned, SDK executes the Python function with the model's arguments"),
        ("Observe",  "Tool result is appended to the Thread as a new message"),
        ("Repeat",   "Loop continues until model returns a final answer (no tool call) or max_turns is hit"),
    ]

    top = Inches(2.05)
    for i, (label, desc) in enumerate(steps, 1):
        _add_text(s, f"{i}  {label}", Inches(0.8), top, Inches(2.0), Inches(0.52),
                  size=16, bold=True, color=GREEN)
        _add_text(s, desc, Inches(2.9), top, Inches(9.6), Inches(0.52),
                  size=15, color=WHITE)
        top += Inches(0.85)

    _add_text(s,
              "Tip: Runner.run() handles the loop for you — but knowing it runs helps you debug.",
              Inches(0.8), Inches(6.38), Inches(11.7), Inches(0.45),
              size=12, bold=True, color=ACCENT, italic=True)
    _add_footer(s)


def slide_06_sdk_core_objects(prs):
    """Concept: SDK core objects."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _section_header(s, "SDK Core Objects",
                    "Four building blocks you'll use in every agent you write")

    objects = [
        ("Agent",   "Declares the agent's identity: instructions (system prompt), model, and list of tools"),
        ("Runner",  "Executes the agent loop synchronously or asynchronously; returns a Result object"),
        ("Tool",    "A Python function decorated with @function_tool; the SDK generates the JSON schema automatically"),
        ("Thread",  "Stateful message history — persisted between runs to give agents memory across turns"),
    ]

    top = Inches(2.05)
    for name, desc in objects:
        _add_text(s, name, Inches(0.8), top, Inches(2.1), Inches(0.55),
                  size=17, bold=True, color=BLUE)
        _add_text(s, desc, Inches(3.0), top, Inches(9.5), Inches(0.55),
                  size=15, color=WHITE)
        top += Inches(1.05)

    _add_footer(s)


def slide_07_tool_definition(prs):
    """Concept: Defining tools."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _section_header(s, "Tool Definition",
                    "How to expose Python functions to your agent")

    points = [
        "Decorate any Python function with @function_tool — no additional config needed",
        "Use type hints and a docstring: the SDK converts them into the JSON schema for the API",
        "Parameters become the tool's input schema; return value (str or Pydantic model) is the output",
        "Raise exceptions or return error strings — the agent will observe failures and can retry or pivot",
        "Tools run in the same process by default; async tools supported with async def",
    ]

    top = Inches(2.2)
    for point in points:
        _add_text(s, f"▸  {point}",
                  Inches(0.8), top, Inches(11.7), Inches(0.88),
                  size=16, color=WHITE)
        top += Inches(0.88)

    _add_text(s,
              'Example:  @function_tool  def get_weather(city: str) -> str: ...',
              Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.45),
              size=13, bold=True, color=GREEN, italic=True)
    _add_footer(s)


def slide_08_runtime_environment(prs):
    """Concept: Runtime environment setup."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _section_header(s, "Runtime Environment Setup",
                    "From zero to a running agent in under 5 minutes")

    steps = [
        ("Requirements", "Python 3.10+  ·  pip or uv  ·  OpenAI API key (platform.openai.com)"),
        ("Install",      "pip install openai-agents   (or: uv add openai-agents)"),
        ("Configure",    "export OPENAI_API_KEY=sk-...   (or set in .env + python-dotenv)"),
        ("First run",    "from agents import Agent, Runner  →  result = Runner.run_sync(agent, 'Hello')"),
        ("Trace viewer", "openai.com/trace — visual replay of every agent loop step; invaluable for debugging"),
    ]

    top = Inches(2.0)
    for label, detail in steps:
        _add_text(s, label, Inches(0.8), top, Inches(2.3), Inches(0.55),
                  size=15, bold=True, color=YELLOW)
        _add_text(s, detail, Inches(3.2), top, Inches(9.3), Inches(0.55),
                  size=15, color=WHITE)
        top += Inches(0.88)

    _add_footer(s)


def slide_09_exercise(prs):
    """Hands-on exercise: Hello World agent."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _section_header(s, "Hands-on Exercise",
                    "Build your first agent end-to-end  — est. 15 min")

    steps = [
        "Install the SDK:  pip install openai-agents  and set OPENAI_API_KEY in your environment",
        "Define one tool:  @function_tool  def echo(message: str) -> str: return f'Echo: {message}'",
        "Create an Agent:  agent = Agent(name='hello', instructions='You are a helpful assistant.', tools=[echo])",
        "Run it:  result = Runner.run_sync(agent, 'Say hello back to me using the echo tool')",
        "Inspect the result:  print(result.final_output)  and explore  result.new_messages  for the loop trace",
    ]

    top = Inches(2.15)
    for i, step in enumerate(steps, 1):
        _add_text(s, f"{i}.", Inches(0.8), top, Inches(0.45), Inches(0.5),
                  size=16, bold=True, color=BLUE)
        _add_text(s, step, Inches(1.3), top, Inches(11.2), Inches(0.75),
                  size=14, color=WHITE)
        top += Inches(0.84)

    _add_text(s, "✓  Done when: agent calls echo(), receives the result, and prints a final answer",
              Inches(0.8), Inches(6.38), Inches(11.7), Inches(0.45),
              size=12, bold=True, color=GREEN)
    _add_footer(s)


def slide_10_knowledge_check(prs):
    """Knowledge check."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _section_header(s, "Knowledge Check",
                    "Test your understanding before moving to Chapter 2")

    q = (
        "A support bot answers one question per turn, has no tools, and forgets each conversation.\n"
        "A colleague's agent books flights, updates a CRM, and picks up where it left off.\n"
        "Which SDK object keeps state between the agent's runs?"
    )
    options = [
        "A.  Agent — because it holds the instructions",
        "B.  Runner — because it executes each turn",
        "C.  Thread  ✓ — persisted message history shared across runs",
        "D.  Tool — because it writes to external systems",
    ]

    _add_text(s, q, Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.3),
              size=15, bold=True, color=WHITE)

    top = Inches(3.4)
    for opt in options:
        color = BLUE if "✓" in opt else WHITE
        _add_text(s, opt, Inches(1.2), top, Inches(11.0), Inches(0.45),
                  size=14, color=color)
        top += Inches(0.5)

    _add_text(s,
              "Answer C: Thread stores the message history. Agent = behavior config. Runner = executor. Tool = action.",
              Inches(0.8), Inches(5.7), Inches(11.7), Inches(0.75),
              size=12, color=ACCENT, italic=True)
    _add_footer(s)


def slide_11_summary(prs):
    """Chapter summary."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _section_header(s, "Chapter 1 Summary",
                    "What you now know")

    points = [
        "Agents differ from chatbots by running a loop, maintaining state, and calling tools autonomously",
        "The agent loop: Receive → Reason → Act → Observe → Repeat (until goal or max_turns)",
        "Four SDK objects: Agent (config), Runner (executor), Tool (function), Thread (memory)",
        "Tools are plain Python functions decorated with @function_tool — SDK handles the schema",
        "You can set up and run a working agent in under 5 minutes with pip install openai-agents",
    ]
    _bullets(s, points, start_top=Inches(2.3), size=16)
    _add_footer(s)


def slide_12_cta(prs):
    """Try it next — CTA slide."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _add_text(s, "Try it next →",
              Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.0),
              size=38, bold=True, align=PP_ALIGN.CENTER)
    _add_text(s, "Chapter 2: The Agent SDK Tool-Calling Model",
              Inches(0.8), Inches(2.7), Inches(11.7), Inches(0.8),
              size=26, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    _add_text(s,
              "You can now build a Hello World agent.\n"
              "Chapter 2 teaches you how to define robust, production-grade tools with Pydantic schemas,\n"
              "handle structured output, and recover gracefully from tool errors.",
              Inches(1.0), Inches(3.7), Inches(11.3), Inches(2.1),
              size=17, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(s, "academy.kspl.tech | Koenig AI Academy",
              Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.4),
              size=10, color=ACCENT, align=PP_ALIGN.CENTER)


# ── main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_01_title(prs)
    slide_02_agenda(prs)
    slide_03_objectives(prs)
    slide_04_agents_vs_chatbots(prs)
    slide_05_agent_loop(prs)
    slide_06_sdk_core_objects(prs)
    slide_07_tool_definition(prs)
    slide_08_runtime_environment(prs)
    slide_09_exercise(prs)
    slide_10_knowledge_check(prs)
    slide_11_summary(prs)
    slide_12_cta(prs)

    out = "vault/courses/openai-agents-sdk-mastery/ch01-slides.pptx"
    prs.save(out)
    slide_count = len(prs.slides)
    print(f"Saved {slide_count} slides → {out}")
    return out, slide_count


if __name__ == "__main__":
    main()
