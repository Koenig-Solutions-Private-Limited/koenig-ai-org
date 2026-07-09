"""Build ch05-slides.pptx for Production Agents: Claude Agent SDK + MCP Connector.

Chapter 5: Production — deploy + observability + cost controls
7 slides: title + 5 content + CTA ("Try It Next").

Fixes KOEA-10989 G0 blockers:
  - All hook slides have 3-5 prose bullets (no code-line text boxes)
  - No 140-char truncation — all bullets manually curated, <120 chars each
  - No wikilinks or inline code noise — clean prose throughout
"""
import sys
import os
sys.path.insert(0, '/tmp/pptx_deps')

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'python-pptx',
                           '--quiet', '--target', '/tmp/pptx_deps'])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

# Brand colours
DARK   = RGBColor(0x18, 0x18, 0x1B)   # #18181b dark background
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)   # #ffffff text
ACCENT = RGBColor(0xA1, 0xA1, 0xAA)  # #a1a1aa zinc-400
INDIGO = RGBColor(0x63, 0x66, 0xF1)  # #6366f1 highlight / CTA

W = Inches(13.333)
H = Inches(7.5)
ATTRIBUTION = "academy.kspl.tech | Koenig AI Academy"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]  # fully blank layout


# ── helpers ───────────────────────────────────────────────────────────────────

def _bg(slide):
    s = slide.shapes.add_shape(1, 0, 0, W, H)
    s.fill.solid()
    s.fill.fore_color.rgb = DARK
    s.line.fill.background()


def _rect(slide, left, top, width, height, fill):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def _text(slide, text, left, top, width, height,
          size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _footer(slide):
    _text(slide, ATTRIBUTION,
          Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
          size=10, color=ACCENT)


def _rule(slide, y=Inches(1.3)):
    r = slide.shapes.add_shape(1, Inches(0.5), y, Inches(12.3), Inches(0.03))
    r.fill.solid()
    r.fill.fore_color.rgb = ACCENT
    r.line.fill.background()


def content_slide(title: str, bullets: list):
    """Add a branded content slide with title + bullet list."""
    slide = prs.slides.add_slide(blank)
    _bg(slide)
    _rect(slide, 0, 0, W, Inches(0.06), INDIGO)
    _text(slide, title,
          Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.95),
          size=28, bold=True)
    _rule(slide, Inches(1.25))

    n = len(bullets)
    step  = Inches(0.90) if n >= 5 else Inches(1.05)
    box_h = Inches(0.80) if n >= 5 else Inches(0.95)
    fsize = 15 if n >= 5 else 17
    y = Inches(1.42)
    for bullet in bullets[:5]:
        _text(slide, f"▸  {bullet}",
              Inches(0.55), y, Inches(12.2), box_h,
              size=fsize, color=WHITE)
        y += step

    _footer(slide)
    return slide


# ── Slide 1: Title ────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank)
_bg(s1)
_rect(s1, 0, 0, W, Inches(0.06), INDIGO)
_text(s1, "CHAPTER 5",
      Inches(1), Inches(1.8), Inches(11.3), Inches(0.5),
      size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
_text(s1, "Production: deploy + observability\n+ cost controls",
      Inches(1), Inches(2.4), Inches(11.3), Inches(2.0),
      size=36, bold=True, align=PP_ALIGN.CENTER)
_text(s1, "Production Agents: Claude Agent SDK + MCP Connector — Koenig AI Academy",
      Inches(1), Inches(5.0), Inches(11.3), Inches(0.5),
      size=14, color=ACCENT, align=PP_ALIGN.CENTER)
_footer(s1)

# ── Slide 2: Hook System Overview ─────────────────────────────────────────────
content_slide(
    "The Hook System — four callbacks, full lifecycle",
    [
        "Hooks are synchronous callbacks in ClaudeAgentOptions.hooks; HookMatcher filters by tool name regex",
        "PreToolUse fires before execution — return permissionDecision: deny to block before any side effect",
        "PostToolUse fires after execution — always returns {}; use for audit logging, not for blocking",
        "Python SDK supports tool, prompt, stop, compaction, permission, notification, subagent events",
        "Python lacks SessionStart/SessionEnd — those lifecycle events are TypeScript-SDK only",
    ]
)

# ── Slide 3: Hook 1 — Audit Log (PostToolUse) ────────────────────────────────
content_slide(
    "Hook 1 — Audit Log (PostToolUse)",
    [
        "PostToolUse fires after every matched tool call, receiving tool_name, file_path, and session_id",
        "Always return {} to pass through — PostToolUse cannot block a write that has already executed",
        "When NOT to use: cost circuit breakers — the file is already written when PostToolUse fires",
        "Log structured JSON, not print() — structured logs enable per-session queries and error rate dashboards",
        "Footgun: if your logger raises an exception, wrap it in try/except or audit entries are silently lost",
    ]
)

# ── Slide 4: Hook 2 — Cost Circuit Breaker (PreToolUse) ─────────────────────
content_slide(
    "Hook 2 — Cost Circuit Breaker (PreToolUse)",
    [
        "PreToolUse fires before the tool executes — return permissionDecision: deny to block before side effects",
        "Return {} to allow; return hookSpecificOutput with permissionDecision and a reason string to deny",
        "Track cumulative input tokens in a class instance; trip the breaker when usage exceeds the cap",
        "Production gotcha: never block silently — always include a reason so Claude gets actionable feedback",
        "Sonnet 5 migration: rebaseline your token cap — Sonnet 5 generates ~30% more tokens per equivalent prompt",
    ]
)

# ── Slide 5: Hook 3 — Session Lifecycle Telemetry ────────────────────────────
content_slide(
    "Hook 3 — Session Lifecycle Telemetry",
    [
        "Python SDK has no SessionStart or SessionEnd events — emit session-start from the first message instead",
        "TypeScript SDK: register hooks.SessionStart with an async callback that fires before any tool call",
        "The SessionStart hook receives session_id, cwd, and environment context in the input dict",
        "Python workaround: log from the first SystemMessage to capture session_id before tool calls begin",
        "When NOT to use TypeScript sessions in Python: the callback is simply not invoked; no error is raised",
    ]
)

# ── Slide 6: Hook 4 — Prompt Sanitization (UserPromptSubmit) ─────────────────
content_slide(
    "Hook 4 — Prompt Sanitization (UserPromptSubmit)",
    [
        "UserPromptSubmit fires before the prompt reaches the model — strip PII at this point, not in PostToolUse",
        "Return a modified input_data dict with the cleaned prompt to replace the original message in the pipeline",
        "When to use: PII redaction, profanity filters, prompt-injection detection, context injection",
        "When NOT to use: tool result cleanup — use PostToolUse for that; this hook only sees the initial prompt",
        "Footgun: always log when redaction occurs with session_id and pattern found — silent redaction fails audits",
    ]
)

# ── Slide 7: CTA ──────────────────────────────────────────────────────────────
s7 = prs.slides.add_slide(blank)
_bg(s7)
_rect(s7, 0, 0, Inches(0.06), H, INDIGO)
_text(s7, "Try It Next",
      Inches(1), Inches(2.0), Inches(11.3), Inches(0.8),
      size=36, bold=True, align=PP_ALIGN.CENTER)
_text(s7, "Add the production hook stack to an agent — verify the circuit breaker fires at a low token cap",
      Inches(1), Inches(3.0), Inches(11.3), Inches(0.7),
      size=20, color=ACCENT, align=PP_ALIGN.CENTER)
_text(s7, ATTRIBUTION,
      Inches(1), Inches(4.2), Inches(11.3), Inches(0.5),
      size=16, color=INDIGO, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
OUT = "vault/courses/production-agents-claude-agent-sdk-mcp-connector/ch05-slides.pptx"
prs.save(OUT)
slide_count = len(prs.slides)
size_kb = os.path.getsize(OUT) // 1024
print(f"OK  {OUT}  ({slide_count} slides, {size_kb} KB)")
