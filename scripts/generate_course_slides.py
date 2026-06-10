#!/usr/bin/env python3
"""Generate ch<NN>-slides.pptx for course chapters using python-pptx.

Usage: generate_course_slides.py <chapter.md> <output.pptx>

Parses 'Slide outline for Slide+Audio Producer' section if present,
otherwise falls back to H2 sections (5 max).
Brand: #18181b dark bg / #ffffff white text / #a1a1aa zinc-400 accent
Format: title + 3-5 content slides + "Try it next" CTA (5-7 total)
"""
import sys
sys.path.insert(0, '/tmp/pptx_deps')

import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

try:
    import yaml
    HAS_YAML = True
except ImportError:
    HAS_YAML = False

DARK   = RGBColor(0x18, 0x18, 0x1B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xA1, 0xA1, 0xAA)

ATTRIBUTION = "academy.kspl.tech | Koenig AI Academy"

SKIP_SECTIONS = {
    "references", "further reading", "related content",
    "notes", "acknowledgements", "slide outline for slide+audio producer",
    "what's next",
}


def parse_frontmatter(text: str):
    if not text.startswith("---"):
        return {}, text
    end = text.find("---", 3)
    if end == -1:
        return {}, text
    fm_raw = text[3:end]
    body = text[end + 3:].strip()
    if HAS_YAML:
        try:
            fm = yaml.safe_load(fm_raw) or {}
        except Exception:
            fm = {}
    else:
        fm = {}
        for line in fm_raw.splitlines():
            if ":" in line:
                k, _, v = line.partition(":")
                fm[k.strip()] = v.strip().strip('"')
    return fm, body


def clean_line(line: str) -> str:
    line = re.sub(r'\*\*(.+?)\*\*', r'\1', line)
    line = re.sub(r'\*(.+?)\*', r'\1', line)
    line = re.sub(r'`([^`]+)`', r'\1', line)
    line = re.sub(r'\[\[([^\]|]+)(?:\|[^\]]+)?\]\]', r'\1', line)
    line = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', line)
    line = re.sub(r'<[^>]+>', '', line)
    line = re.sub(r'^\s*[-*+]\s+', '', line)
    line = re.sub(r'^\s*\d+\.\s+', '', line)
    line = re.sub(r'\[\^[^\]]+\]', '', line)
    return line.strip()


def extract_slide_outline(body: str):
    """Return list of (title, []) from the Slide outline section if present."""
    m = re.search(
        r'## Slide outline for Slide\+Audio Producer\s*\n(.*?)(?=\n## |\Z)',
        body, re.IGNORECASE | re.DOTALL
    )
    if not m:
        return []
    block = m.group(1)
    slides = []
    for raw in block.splitlines():
        cleaned = clean_line(raw)
        # Lines like "Slide N: <title>" or "- Slide N: <title>"
        m2 = re.match(r'^(?:Slide\s+\d+\s*:\s*)(.+)', cleaned, re.IGNORECASE)
        if m2:
            slides.append(m2.group(1).strip())
    return slides


def extract_h2_sections(body: str, skip: set):
    """Return list of (h2_title, [bullet_str, ...])."""
    sections = []
    current_title = None
    bullets: list[str] = []

    skip_lower = {s.lower() for s in skip}

    for raw_line in body.splitlines():
        if raw_line.startswith("## "):
            if current_title is not None and current_title.lower() not in skip_lower:
                sections.append((current_title, bullets[:5]))
            current_title = raw_line[3:].strip()
            bullets = []
        elif raw_line.startswith("### "):
            sub = raw_line[4:].strip()
            if sub:
                bullets.append(sub[:100])
        elif current_title is not None:
            cleaned = clean_line(raw_line)
            if (cleaned
                    and len(cleaned) > 15
                    and not cleaned.startswith("|")
                    and not cleaned.startswith("```")
                    and not cleaned.startswith("![")):
                bullets.append(cleaned[:120])

    if current_title is not None and current_title.lower() not in skip_lower:
        sections.append((current_title, bullets[:5]))

    return [(t, b) for t, b in sections if b]


def bullets_for_slide(title: str, h2_sections: list) -> list[str]:
    """Find the best matching H2 section bullets for an outline slide title."""
    title_lower = title.lower()
    # Score each section by keyword overlap
    best_score = -1
    best_bullets: list[str] = []
    title_words = set(re.findall(r'\w+', title_lower))

    for sec_title, sec_bullets in h2_sections:
        sec_words = set(re.findall(r'\w+', sec_title.lower()))
        score = len(title_words & sec_words)
        if score > best_score and sec_bullets:
            best_score = score
            best_bullets = sec_bullets

    return best_bullets[:4] if best_bullets else []


def _add_text(slide, text, left, top, width, height,
              size, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _set_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK


def _add_footer(slide):
    _add_text(slide, ATTRIBUTION,
              Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.4),
              size=10, color=ACCENT)


def make_title_slide(prs, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, title,
              Inches(1.0), Inches(2.0), Inches(11.3), Inches(2.5),
              size=34, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        _add_text(slide, subtitle,
                  Inches(1.0), Inches(4.7), Inches(11.3), Inches(0.5),
                  size=16, color=ACCENT, align=PP_ALIGN.CENTER)
    _add_footer(slide)


def make_content_slide(prs, section_title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_text(slide, section_title,
              Inches(0.8), Inches(0.45), Inches(11.7), Inches(1.1),
              size=26, bold=True)
    rect = slide.shapes.add_shape(
        1,
        Inches(0.8), Inches(1.55), Inches(11.7), Inches(0.04)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = ACCENT
    rect.line.fill.background()

    top = Inches(1.75)
    for bullet in bullets[:4]:
        _add_text(slide, f"  {bullet}",
                  Inches(0.8), top, Inches(11.7), Inches(1.05),
                  size=17, color=WHITE)
        top += Inches(1.15)

    _add_footer(slide)


def make_cta_slide(prs, next_chapter: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    cta_text = f"Try it next{': ' + next_chapter if next_chapter else ''}"
    _add_text(slide, cta_text,
              Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.0),
              size=28, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "academy.kspl.tech",
              Inches(1.0), Inches(3.5), Inches(11.3), Inches(0.9),
              size=36, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    _add_text(slide, "Koenig AI Academy",
              Inches(1.0), Inches(6.85), Inches(11.3), Inches(0.4),
              size=12, color=ACCENT, align=PP_ALIGN.CENTER)


def generate(chapter_path: Path, output_path: Path) -> Path:
    text = chapter_path.read_text(encoding="utf-8")
    fm, body = parse_frontmatter(text)

    title = fm.get("title") or ""
    if not title:
        h1 = re.search(r'^# (.+)', body, re.MULTILINE)
        title = h1.group(1).strip() if h1 else chapter_path.stem.replace("-", " ").title()
    # Strip surrounding quotes if any
    title = title.strip('"\'')

    course_slug = fm.get("course_slug", "")
    chapter_num = fm.get("chapter_num", "")
    subtitle_parts = []
    if chapter_num:
        subtitle_parts.append(f"Chapter {chapter_num}")
    if course_slug:
        subtitle_parts.append(course_slug.replace("-", " ").title())
    subtitle = " · ".join(subtitle_parts)

    # Extract next-chapter hint from "What's next" section
    next_chapter = ""
    m = re.search(r"## What.s next\s*\n(.+?)(?=\n##|\Z)", body, re.DOTALL | re.IGNORECASE)
    if m:
        first_sentence = re.split(r'[.!?]', m.group(1).strip())[0]
        next_chapter = clean_line(first_sentence)[:80]

    # Get H2 sections for bullet extraction
    h2_sections = extract_h2_sections(body, SKIP_SECTIONS)

    # Try to get the embedded slide outline
    outline_titles = extract_slide_outline(body)

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs, title, subtitle)

    if outline_titles:
        # Use outline titles (skip first and last — those are intro/bridge; use middle 5)
        content_titles = outline_titles[1:-1] if len(outline_titles) > 2 else outline_titles
        content_titles = content_titles[:5]  # cap at 5 for 7-slide total
        for slide_title in content_titles:
            bullets = bullets_for_slide(slide_title, h2_sections)
            if not bullets:
                # Generate bullets directly from the title text
                bullets = [slide_title]
            make_content_slide(prs, slide_title, bullets)
    else:
        # Fallback: use H2 sections
        for sec_title, bullets in h2_sections[:5]:
            make_content_slide(prs, sec_title, bullets)

    make_cta_slide(prs, next_chapter)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: generate_course_slides.py <chapter.md> <output.pptx>")
        sys.exit(1)

    chapter_p = Path(sys.argv[1])
    output_p  = Path(sys.argv[2])

    if not chapter_p.exists():
        print(f"ERROR: chapter not found: {chapter_p}")
        sys.exit(1)

    try:
        out = generate(chapter_p, output_p)
        prs = Presentation(str(out))
        slide_count = len(prs.slides)
        size_kb = out.stat().st_size // 1024
        print(f"OK  {out}  ({slide_count} slides, {size_kb} KB)")
    except Exception as e:
        import traceback
        traceback.print_exc()
        print(f"FAIL: {e}")
        sys.exit(1)
