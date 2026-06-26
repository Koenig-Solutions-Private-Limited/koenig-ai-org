"""Build ch02-slides.pptx for MCP: From First Principles to Production — Chapter 2.

Chapter: JSON-RPC over stdio — the wire protocol explained
7 slides: title + 5 content + CTA ("Try It Next").
Fixes KOEA-8271 G0 blockers:
  - No truncated text boxes (all bullets manually curated, <100 chars each)
  - No wikilink / footnote notation — live URLs used where references cited
"""
import sys
import os
sys.path.insert(0, '/tmp/pptx_deps')

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'python-pptx',
                           '--quiet', '--target', '/tmp/pptx_deps'])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

# Brand colours
DARK   = RGBColor(0x18, 0x18, 0x1B)   # #18181b dark background
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)   # #ffffff text
ACCENT = RGBColor(0xA1, 0xA1, 0xAA)  # #a1a1aa zinc-400
INDIGO = RGBColor(0x63, 0x66, 0xF1)  # #6366f1 highlight / CTA

W = Inches(13.333)
H = Inches(7.5)
ATTRIBUTION = "academy.kspl.tech | Koenig AI Academy"

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]  # fully blank layout


# ── helpers ───────────────────────────────────────────────────────────────────

def _bg(slide):
    s = slide.shapes.add_shape(1, 0, 0, W, H)
    s.fill.solid()
    s.fill.fore_color.rgb = DARK
    s.line.fill.background()


def _rect(slide, left, top, width, height, fill):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def _text(slide, text, left, top, width, height,
          size=16, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _footer(slide):
    _text(slide, ATTRIBUTION,
          Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
          size=10, color=ACCENT)


def _rule(slide, y=Inches(1.3)):
    r = slide.shapes.add_shape(1, Inches(0.5), y, Inches(12.3), Inches(0.03))
    r.fill.solid()
    r.fill.fore_color.rgb = ACCENT
    r.line.fill.background()


def content_slide(title: str, bullets: list):
    """Add a branded content slide with title + bullet list."""
    slide = prs.slides.add_slide(blank)
    _bg(slide)
    _rect(slide, 0, 0, W, Inches(0.06), INDIGO)
    _text(slide, title,
          Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.95),
          size=28, bold=True)
    _rule(slide, Inches(1.25))

    n = len(bullets)
    step  = Inches(0.90) if n >= 5 else Inches(1.05)
    box_h = Inches(0.80) if n >= 5 else Inches(0.95)
    fsize = 15 if n >= 5 else 17
    y = Inches(1.42)
    for bullet in bullets[:5]:
        _text(slide, f"▸  {bullet}",
              Inches(0.55), y, Inches(12.2), box_h,
              size=fsize, color=WHITE)
        y += step

    _footer(slide)
    return slide


# ── Slide 1: Title ────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(blank)
_bg(s1)
_rect(s1, 0, 0, W, Inches(0.06), INDIGO)
_text(s1, "CHAPTER 2",
      Inches(1), Inches(1.8), Inches(11.3), Inches(0.5),
      size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
_text(s1, "JSON-RPC over stdio\n— the wire protocol explained",
      Inches(1), Inches(2.4), Inches(11.3), Inches(2.0),
      size=36, bold=True, align=PP_ALIGN.CENTER)
_text(s1, "MCP: From First Principles to Production — Koenig AI Academy",
      Inches(1), Inches(5.0), Inches(11.3), Inches(0.5),
      size=14, color=ACCENT, align=PP_ALIGN.CENTER)
_footer(s1)

# ── Slide 2: Why the wire format matters ──────────────────────────────────────
content_slide(
    "Why the wire format matters",
    [
        "All MCP SDKs abstract over the same wire protocol — production failures reduce to raw JSON",
        "Every SDK call maps to one or two JSON messages; wire fluency exposes round-trip costs",
        "Wire format knowledge reveals error shapes, security exposure, and protocol invariants",
        "SDK fluency without wire fluency: you can configure MCP integrations but not diagnose them",
    ]
)

# ── Slide 3: JSON-RPC 2.0 — the message envelope ──────────────────────────────
content_slide(
    "JSON-RPC 2.0 — the message envelope",
    [
        "MCP uses JSON-RPC 2.0 as its message format — a deliberate choice, not an implementation detail",
        "The id field correlates responses to requests; enables multiple in-flight requests simultaneously",
        "Notifications have no id field — one-way messages; the sender never expects a response",
        "error and result are mutually exclusive in a response — a valid response carries exactly one",
        "Spec: https://www.jsonrpc.org/specification",
    ]
)

# ── Slide 4: stdio transport ──────────────────────────────────────────────────
content_slide(
    "stdio transport — why a pipe beats a socket",
    [
        "Host launches MCP server as a subprocess; communication is via stdin and stdout pipes",
        "Framing rule: read until newline, parse as JSON — no sockets, no ports, no TLS needed",
        "Subprocess lifecycle: server dies when host closes the pipe — no cleanup logic required",
        "stderr is strictly for logs — any non-JSON text on stdout corrupts the JSON-RPC stream",
        "Spec: https://spec.modelcontextprotocol.io/specification/2025-03-26/basic/transports/",
    ]
)

# ── Slide 5: The initialize handshake ─────────────────────────────────────────
content_slide(
    "The initialize handshake",
    [
        "Every session begins: initialize request → initialize response → notifications/initialized",
        "The handshake must complete before any tool calls — requests before it violate the spec",
        "Each side’s capabilities object determines the entire valid call surface for the session",
        "notifications/initialized is a notification — no id field, no response expected from server",
        "Spec: https://spec.modelcontextprotocol.io/specification/2025-03-26/",
    ]
)

# ── Slide 6: Reading a real MCP exchange ──────────────────────────────────────
content_slide(
    "Reading a real MCP exchange",
    [
        "Scenario: user asks about open GitHub PRs; host invokes list_pull_requests on an MCP server",
        "Step 1: host sends tools/list to discover available tools and each tool’s inputSchema",
        "Step 2: model decides to call a tool; host sends tools/call with name and arguments",
        "Tool errors use isError:true in result — JSON-RPC error field is for protocol errors only",
        "Protocol errors (malformed request) and tool errors (execution failed) are handled differently",
    ]
)

# ── Slide 7: CTA ──────────────────────────────────────────────────────────────
s7 = prs.slides.add_slide(blank)
_bg(s7)
_rect(s7, 0, 0, Inches(0.06), H, INDIGO)
_text(s7, "Try It Next",
      Inches(1), Inches(2.0), Inches(11.3), Inches(0.8),
      size=36, bold=True, align=PP_ALIGN.CENTER)
_text(s7, "Write a 60-line Python MCP server — no SDK, raw sys.stdin / sys.stdout",
      Inches(1), Inches(3.0), Inches(11.3), Inches(0.7),
      size=20, color=ACCENT, align=PP_ALIGN.CENTER)
_text(s7, ATTRIBUTION,
      Inches(1), Inches(4.2), Inches(11.3), Inches(0.5),
      size=16, color=INDIGO, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
OUT = "vault/courses/mcp-from-first-principles-to-production/ch02-slides.pptx"
prs.save(OUT)
slide_count = len(prs.slides)
size_kb = os.path.getsize(OUT) // 1024
print(f"OK  {OUT}  ({slide_count} slides, {size_kb} KB)")
