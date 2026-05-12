#!/usr/bin/env python3
"""Generate ch03-slides.pptx for picking-a-frontier-model-2026-q2.

Matches the visual style of ch01, ch02, ch04 decks:
- 13.33×7.5 in slides
- Brand: #18181b dark bg / #ffffff white / #a1a1aa zinc-400 accent / #3b82f6 blue highlight
- Structure: chapter header · learning objectives · content slides · CTA
"""
import sys
sys.path.insert(0, '/tmp/pptx_deps')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK   = RGBColor(0x18, 0x18, 0x1B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xA1, 0xA1, 0xAA)   # zinc-400
BLUE   = RGBColor(0x3B, 0x82, 0xF6)   # blue-500
YELLOW = RGBColor(0xFA, 0xCC, 0x15)   # yellow-400

COURSE = "Picking a Frontier Model — Q2 2026"
FOOTER = "academy.kspl.tech | Koenig AI Academy"


# ── helpers ───────────────────────────────────────────────────────────────────

def _set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = DARK


def _add_text(slide, text, left, top, width, height,
              size, bold=False, color=WHITE,
              align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _add_footer(slide, text=FOOTER):
    _add_text(slide, text,
              Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.4),
              size=10, color=ACCENT)


def _add_rule(slide, top=Inches(1.55)):
    rect = slide.shapes.add_shape(1, Inches(0.8), top, Inches(11.7), Inches(0.04))
    rect.fill.solid()
    rect.fill.fore_color.rgb = ACCENT
    rect.line.fill.background()


def _add_section_header(slide, title, subtitle=None):
    """Big title + optional subtitle for content slides."""
    _add_text(slide, title,
              Inches(0.8), Inches(0.35), Inches(11.7), Inches(1.1),
              size=28, bold=True)
    _add_rule(slide)
    if subtitle:
        _add_text(slide, subtitle,
                  Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.6),
                  size=14, color=ACCENT, italic=True)


def _bullets(slide, items, start_top=None, size=17, indent="▸  "):
    top = start_top if start_top is not None else Inches(2.3)
    for item in items:
        _add_text(slide, f"{indent}{item}",
                  Inches(0.8), top, Inches(11.7), Inches(1.0),
                  size=size, color=WHITE)
        top += Inches(0.92)


# ── slide builders ────────────────────────────────────────────────────────────

def slide_01_title(prs):
    """Chapter header — full-bleed dark."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _add_text(s, "PICKING A FRONTIER MODEL · 2026 Q2",
              Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.6),
              size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    _add_text(s, "Chapter 3",
              Inches(0.8), Inches(1.85), Inches(11.7), Inches(0.7),
              size=22, bold=False, color=ACCENT, align=PP_ALIGN.CENTER)
    _add_text(s, "Long-context behavior\n— effective vs. advertised context windows",
              Inches(0.8), Inches(2.7), Inches(11.7), Inches(2.2),
              size=36, bold=True, align=PP_ALIGN.CENTER)
    _add_text(s, "academy.kspl.tech | Koenig AI Academy",
              Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.4),
              size=10, color=ACCENT, align=PP_ALIGN.CENTER)


def slide_02_objectives(prs):
    """Learning objectives."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _add_section_header(s, "What you'll be able to do", "Learning Objectives")
    objectives = [
        "Define effective context window vs. advertised limit — why they differ",
        "Run a needle-in-haystack test at 50K / 200K / 500K token depths",
        "Identify the 3 failure modes at scale: lost needles · hallucinated synthesis · degraded reasoning",
        "Choose the right context strategy — full context vs. RAG — for your document volume",
        "Understand why 1M token context ≠ 1M token understanding",
    ]
    _bullets(s, objectives, start_top=Inches(2.3), size=16)
    _add_footer(s)


def slide_03_key_facts(prs):
    """7 key facts — the empirical baseline."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _add_section_header(s, "Key Facts", "The empirical baseline before you test")
    facts = [
        "Lost-in-the-middle: 20–40 pp accuracy drop for info buried in middle of long context  [Liu et al. 2023]",
        "Gemini 3.1 Pro leads raw single-fact retrieval up to ~600K tokens",
        "Opus 4.7 synthesis effective limit (~500K) beats Gemini's (~300K) on multi-hop tasks",
        "GPT-5.5 (128K window) shows the most stable middle-position retrieval of the three",
        "Three-or-more-fact synthesis is unreliable beyond 200K tokens on ALL three models",
        "RAG (top-k=5) is 10–20× cheaper and matches full-context synthesis accuracy for docs >100K tokens",
    ]
    _bullets(s, facts, start_top=Inches(2.2), size=15)
    _add_footer(s)


def slide_04_effective_window(prs):
    """Advertised vs. effective context window with the comparison table."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _add_section_header(s, "Advertised vs. Effective Context Window",
                        "Vendors advertise size. They don't advertise the accuracy curve inside it.")

    # Three concepts
    concepts = [
        ("Retrieval effective limit", "depth where single-fact accuracy < 90%  — safest boundary for fact-lookup"),
        ("Synthesis effective limit", "depth where cross-doc reasoning < 85%  — typically 30-50% of retrieval limit"),
        ("Hot zone",                  "first & last ~15% of context — all models show dramatically higher accuracy here"),
    ]
    top = Inches(2.25)
    for term, desc in concepts:
        _add_text(s, term, Inches(0.8), top, Inches(3.8), Inches(0.55),
                  size=15, bold=True, color=BLUE)
        _add_text(s, desc, Inches(4.7), top, Inches(7.8), Inches(0.55),
                  size=15, color=WHITE)
        top += Inches(0.62)

    # Comparison table header
    _add_rule(s, top=Inches(4.2))
    top = Inches(4.3)
    cols = ["Model", "Advertised", "Retrieval effective", "Synthesis effective"]
    widths = [2.0, 2.0, 3.0, 3.0]
    left = 0.8
    for col, w in zip(cols, widths):
        _add_text(s, col, Inches(left), top, Inches(w), Inches(0.45),
                  size=13, bold=True, color=ACCENT)
        left += w

    rows = [
        ["Opus 4.7",       "1M",    "~800K", "~500K"],
        ["GPT-5.5",        "128K",  "~120K", "~75K"],
        ["Gemini 3.1 Pro", "1M",    "~700K", "~300K"],
    ]
    top += Inches(0.48)
    for row in rows:
        left = 0.8
        for cell, w in zip(row, widths):
            _add_text(s, cell, Inches(left), top, Inches(w), Inches(0.45),
                      size=14, color=WHITE)
            left += w
        top += Inches(0.5)

    _add_footer(s)


def slide_05_failure_modes(prs):
    """Three failure modes at scale."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _add_section_header(s, "The 3 Failure Modes at Scale",
                        "Recognise them before they reach production")

    modes = [
        ("1  Lost Needles (retrieval miss)",
         "Model ignores a fact that IS in context — not hallucinated, simply not retrieved\n"
         "Most common at 50K–200K (GPT-5.5) or 200K–500K (Gemini 3.1 Pro)"),
        ("2  Hallucinated Synthesis",
         "Model combines real retrieved facts with invented connections — fluent, confident, partially fabricated\n"
         "Harder to detect — requires ground-truth verification"),
        ("3  Degraded Step-by-Step Reasoning",
         "Chain-of-thought shortens at high depth; model skips intermediate steps it handles at low depth\n"
         "Appears in math word problems, multi-step code analysis, legal document reasoning"),
    ]
    top = Inches(2.1)
    for title, body in modes:
        _add_text(s, title, Inches(0.8), top, Inches(11.7), Inches(0.5),
                  size=16, bold=True, color=YELLOW)
        top += Inches(0.52)
        _add_text(s, body, Inches(1.2), top, Inches(11.3), Inches(0.85),
                  size=13, color=WHITE)
        top += Inches(0.98)

    _add_footer(s)


def slide_06_needle_haystack(prs):
    """Needle-in-haystack evaluation methodology."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _add_section_header(s, "The Needle-in-Haystack Evaluation",
                        "Standard method for measuring retrieval effective limit")

    steps = [
        "Prepare a haystack — large document padded to target token depth (legal corpus, Wikipedia, synthetic filler)",
        'Insert a needle — unique fact that cannot be guessed: "The secret phrase is: banana-lighthouse-44"',
        "Insert needle at a specific depth (e.g. 25% / 50% / 75% of total context)",
        "Ask the model to retrieve the needle  →  correct = 1, any other answer = 0",
        "Repeat across a grid:  context size (50K / 100K / 200K / 500K)  ×  position (10% / 25% / 50% / 75% / 90%)",
        "≥3 runs per cell to average noise  →  output an accuracy heatmap",
    ]
    top = Inches(2.1)
    for i, step in enumerate(steps, 1):
        _add_text(s, f"{i}.", Inches(0.8), top, Inches(0.45), Inches(0.5),
                  size=16, bold=True, color=BLUE)
        _add_text(s, step, Inches(1.3), top, Inches(11.2), Inches(0.5),
                  size=15, color=WHITE)
        top += Inches(0.73)

    _add_footer(s)


def slide_07_context_strategy(prs):
    """Choosing your context strategy — decision table."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _add_section_header(s, "Choosing Your Context Strategy",
                        "Practical decision framework for multi-document workloads")

    rows = [
        ("< 50K tokens",         "Full context (any model)",                        "All three models reliable below 50K"),
        ("50K – 120K",           "Full context; test empirically",                  "GPT-5.5 shows good middle-position stability here"),
        ("120K – 500K",          "Opus 4.7 full context  OR  RAG pipeline",         "Within Opus 4.7's synthesis effective limit (~500K)"),
        ("500K – 800K",          "Gemini for retrieval; chunked Opus for synthesis", "Both approach synthesis limits; chunking reduces depth"),
        ("> 800K tokens",        "RAG pipeline + any model",                        "Beyond all models' reliable retrieval limits"),
    ]
    headers = ["Document volume", "Strategy", "Rationale"]
    widths  = [2.3, 4.3, 5.1]
    top = Inches(1.85)
    left_base = 0.8

    # Header row
    left = left_base
    for h, w in zip(headers, widths):
        _add_text(s, h, Inches(left), top, Inches(w), Inches(0.45),
                  size=13, bold=True, color=ACCENT)
        left += w
    top += Inches(0.5)

    for row in rows:
        left = left_base
        for cell, w in zip(row, widths):
            _add_text(s, cell, Inches(left), top, Inches(w), Inches(0.55),
                      size=12, color=WHITE)
            left += w
        top += Inches(0.65)

    _add_text(s,
              "Key principle: use long context for retrieval tasks · use chunking + multiple calls for synthesis tasks",
              Inches(0.8), Inches(6.45), Inches(11.7), Inches(0.45),
              size=13, bold=True, color=BLUE)
    _add_footer(s)


def slide_08_rag_vs_fullctx(prs):
    """RAG vs. full-context tradeoff quantified."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _add_section_header(s, "RAG vs. Full-Context — 200K Corpus",
                        "Rough cost + accuracy comparison from internal workloads")

    rows = [
        ("Gemini 3.1 Pro, full ctx",    "$$$",  "94%",  "81%"),
        ("Opus 4.7, full ctx",          "$$",   "91%",  "88%"),
        ("RAG top-k=5 + Opus 4.7",      "$",    "87%*", "92%"),
        ("RAG top-k=5 + GPT-5.5",       "$",    "87%*", "89%"),
    ]
    headers = ["Approach", "Cost", "Retrieval acc.", "Synthesis acc."]
    widths  = [4.2, 1.5, 2.3, 2.3]
    left_base = 0.8
    top = Inches(1.85)

    left = left_base
    for h, w in zip(headers, widths):
        _add_text(s, h, Inches(left), top, Inches(w), Inches(0.45),
                  size=13, bold=True, color=ACCENT)
        left += w
    top += Inches(0.52)

    for row in rows:
        left = left_base
        for cell, w in zip(row, widths):
            _add_text(s, cell, Inches(left), top, Inches(w), Inches(0.52),
                      size=14, color=WHITE)
            left += w
        top += Inches(0.6)

    _add_text(s, "* RAG retrieval limited by embedding step — may miss right chunk",
              Inches(0.8), Inches(5.1), Inches(11.7), Inches(0.4),
              size=11, color=ACCENT, italic=True)

    takeaways = [
        "If workload is primarily SYNTHESIS → use RAG  (10–20× cheaper, matches or beats full-context accuracy)",
        "If workload is primarily SINGLE-FACT RETRIEVAL from one doc → full context wins on simplicity + reliability",
        "For pure retrieval depth, Gemini 3.1 Pro has a genuine advantage over Opus 4.7",
    ]
    top = Inches(5.6)
    for t in takeaways:
        _add_text(s, f"▸  {t}", Inches(0.8), top, Inches(11.7), Inches(0.5),
                  size=13, color=WHITE)
        top += Inches(0.5)

    _add_footer(s)


def slide_09_handson(prs):
    """Hands-on exercise."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _add_section_header(s, "Hands-on Exercise",
                        "Run needle-in-haystack at 3 depth levels on your own documents  — est. 25 min")

    steps = [
        "Choose a document/set from your production context. Prepare versions at ~50K, ~200K, and target depth.",
        "Insert 3 needles: Needle A near start (5–10%) · Needle B in middle (45–55%) · Needle C near end (85–95%)",
        'For each model, ask: "What is the value of [needle identifier]?" — run each ≥3 times',
        "Record a 3×3 accuracy grid (3 depths × 3 positions). Note failure positions and depths.",
        "Run ≥1 multi-hop synthesis task: question requiring facts from both Needles A and C.",
    ]
    top = Inches(2.15)
    for i, step in enumerate(steps, 1):
        _add_text(s, f"{i}.", Inches(0.8), top, Inches(0.45), Inches(0.5),
                  size=16, bold=True, color=YELLOW)
        _add_text(s, step, Inches(1.3), top, Inches(11.2), Inches(0.5),
                  size=15, color=WHITE)
        top += Inches(0.72)

    _add_text(s, "✓  Done when: 3×3 retrieval grid filled · effective limit estimated · multi-hop pass/fail recorded",
              Inches(0.8), Inches(6.42), Inches(11.7), Inches(0.45),
              size=12, bold=True, color=BLUE)
    _add_footer(s)


def slide_10_knowledge_check(prs):
    """Knowledge check teaser."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _add_section_header(s, "Knowledge Check",
                        "Test your understanding before moving to Chapter 4")

    q1 = (
        "A legal tech team needs to extract all liability clauses from 400K–700K-token contracts\n"
        "and synthesise maximum exposure. Which approach is best?"
    )
    options1 = [
        "A.  Gemini 3.1 Pro full context — 1M window covers it",
        "B.  Opus 4.7 full context — best synthesis at any depth",
        "C.  RAG (chunk by clause, embed, retrieve top-k) + Opus 4.7  ✓",
        "D.  GPT-5.5 — most stable middle-position retrieval",
    ]
    _add_text(s, q1, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.9),
              size=15, bold=True, color=WHITE)
    top = Inches(3.05)
    for opt in options1:
        color = BLUE if "✓" in opt else WHITE
        _add_text(s, opt, Inches(1.2), top, Inches(11.0), Inches(0.45),
                  size=14, color=color)
        top += Inches(0.5)

    _add_text(s,
              "Answer C: RAG + Opus 4.7 reduces context depth below Opus's synthesis limit; "
              "clause-level chunking gives a short, high-quality context to reason over.",
              Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.75),
              size=12, color=ACCENT, italic=True)
    _add_footer(s)


def slide_11_cta(prs):
    """Try it next — CTA slide."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s)
    _add_text(s, "Try it next →",
              Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.0),
              size=38, bold=True, align=PP_ALIGN.CENTER)
    _add_text(s, "Chapter 4: Cost-per-task",
              Inches(0.8), Inches(2.7), Inches(11.7), Inches(0.8),
              size=26, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    _add_text(s,
              "You now know which model wins on retrieval depth and which wins on synthesis.\n"
              "Chapter 4 answers the final question: what does reliable output actually cost?\n"
              "Build a cost-per-task model accounting for retry rates, context caching, and tool-call overhead.",
              Inches(1.0), Inches(3.7), Inches(11.3), Inches(2.1),
              size=17, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(s, "academy.kspl.tech | Koenig AI Academy",
              Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.4),
              size=10, color=ACCENT, align=PP_ALIGN.CENTER)


# ── main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_01_title(prs)
    slide_02_objectives(prs)
    slide_03_key_facts(prs)
    slide_04_effective_window(prs)
    slide_05_failure_modes(prs)
    slide_06_needle_haystack(prs)
    slide_07_context_strategy(prs)
    slide_08_rag_vs_fullctx(prs)
    slide_09_handson(prs)
    slide_10_knowledge_check(prs)
    slide_11_cta(prs)

    out = "vault/courses/picking-a-frontier-model-2026-q2/ch03-slides.pptx"
    prs.save(out)
    print(f"Saved {len(prs.slides)} slides → {out}")
    return out


if __name__ == "__main__":
    main()
